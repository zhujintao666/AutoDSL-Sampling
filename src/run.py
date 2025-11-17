# scripts/render_table_from_json.py
# -*- coding: utf-8 -*-
import os, json
from collections import OrderedDict
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(__file__))
OUT = os.path.join(ROOT, "outputs")

# 默认文件路径
P_STRUCTURE = os.path.join(OUT, "final_structure.pruned.json")
P_STATES    = os.path.join(OUT, "final_state_space.json")
P_ACTIONS   = os.path.join(OUT, "final_action_space.json")
P_PRIORITY  = os.path.join(OUT, "final_priority.json")  # 可选

AGENTS = ["Supplier","Manufacturer","Distributor","Wholesaler","Retailer"]

def _load(path, default):
    try:
        with open(path,"r",encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return default

def _lines(xs):
    if not xs: return "—"
    return "\n".join(xs)

def _fmt_objectives(pri):
    """
    兼容几种可能结构：
    1) {agent: [{"Prop": "Inventory","Direction":"Max"}...]}
    2) {agent: [["Inventory","Max"], ...]} 或 {agent: [["Inventory","Increase"], ...]}
    3) {agent: ["Max Inventory","Min Cost", ...]}
    """
    out = {}
    for ag in AGENTS:
        v = pri.get(ag, [])
        items = []
        for it in v:
            if isinstance(it, dict):
                p = it.get("Prop") or it.get("prop") or it.get("name") or it.get("Property")
                d = (it.get("Direction") or it.get("direction") or "").lower()
            elif isinstance(it, (list, tuple)) and len(it)>=2:
                p, d = it[0], str(it[1]).lower()
            else:
                # 直接字符串
                s = str(it)
                items.append(s)
                continue
            if not p: continue
            if d.startswith("max") or "increase" in d:
                items.append(f"Maximize {p}")
            elif d.startswith("min") or "decrease" in d or "reduce" in d:
                items.append(f"Minimize {p}")
            else:
                items.append(f"Optimize {p}")
        out[ag] = items
    return out

def load_all():
    struct  = _load(P_STRUCTURE, {})
    states  = _load(P_STATES,   {})
    actions = _load(P_ACTIONS,  {})
    priority= _load(P_PRIORITY, {})
    if priority:
        objectives = _fmt_objectives(priority)
    else:
        objectives = {ag: [] for ag in AGENTS}

    # 保障 key 顺序 & 缺省
    props = {ag: struct.get(ag, []) for ag in AGENTS}
    st    = {ag: states.get(ag, {"Self":[], "Neighbors":[], "Relations":[]}) for ag in AGENTS}
    acts  = {ag: actions.get(ag, []) for ag in AGENTS}
    return props, st, acts, objectives

def render_html(props, st, acts, objectives, outpath):
    def cell_actions(lst):
        if not lst: return "—"
        lines = []
        for a in lst:
            t = a.get("type","")
            ps = a.get("params",[])
            if ps:
                lines.append(f"{t}: {', '.join(ps)}")
            else:
                lines.append(f"{t}")
        return "<br/>".join(lines)

    css = """
    <style>
    body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, 'PingFang SC', 'Microsoft YaHei', Arial; }
    table { border-collapse: collapse; width: 100%; }
    th, td { border:1px solid #ccc; padding:10px; vertical-align: top; }
    th { background:#f7f7f7; }
    .rowTitle { width:140px; font-weight:700; color:#c62828; }
    .sub { color:#555; font-weight:700; margin-top:6px; }
    .cell p { margin:6px 0; }
    </style>
    """
    html = [f"<!doctype html><meta charset='utf-8'>{css}<table>"]
    # header
    html.append("<tr><th></th>" + "".join([f"<th>{ag}</th>" for ag in AGENTS]) + "</tr>")

    # Properties
    html.append("<tr><td class='rowTitle'>Properties</td>" +
                "".join([f"<td class='cell'><p>{'<br/>'.join(props[ag]) or '—'}</p></td>" for ag in AGENTS]) +
                "</tr>")

    # Objectives
    html.append("<tr><td class='rowTitle'>Objectives</td>" +
                "".join([f"<td class='cell'><p>{'<br/>'.join(objectives.get(ag,[])) or '—'}</p></td>" for ag in AGENTS]) +
                "</tr>")

    # States (Self / Neighbors / Relations)
    def state_cell(ag):
        s = st[ag]
        return (
            f"<div class='sub'>Self:</div><div>{'<br/>'.join(s.get('Self',[])) or '—'}</div>"
            f"<div class='sub'>Neighbors:</div><div>{'<br/>'.join(s.get('Neighbors',[])) or '—'}</div>"
            f"<div class='sub'>Relations:</div><div>{'<br/>'.join(s.get('Relations',[])) or '—'}</div>"
        )
    html.append("<tr><td class='rowTitle'>States</td>" +
                "".join([f"<td class='cell'>{state_cell(ag)}</td>" for ag in AGENTS]) +
                "</tr>")

    # Actions
    html.append("<tr><td class='rowTitle'>Actions</td>" +
                "".join([f"<td class='cell'>{cell_actions(acts[ag])}</td>" for ag in AGENTS]) +
                "</tr>")

    html.append("</table>")
    with open(outpath,"w",encoding="utf-8") as f:
        f.write("\n".join(html))

def render_pptx(props, st, acts, objectives, outpath):
    prs = Presentation()
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    rows = 1 + 4  # header + 4 blocks
    cols = 1 + len(AGENTS)

    left, top, width, height = Inches(0.3), Inches(0.3), Inches(12.9), Inches(6.6)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # widths
    table.columns[0].width = Inches(1.6)
    for i in range(1, cols):
        table.columns[i].width = Inches(2.25)

    # header
    table.cell(0,0).text = ""
    for i,ag in enumerate(AGENTS, start=1):
        table.cell(0,i).text = ag

    def set_cell(r,c,text):
        table.cell(r,c).text = text
        for p in table.cell(r,c).text_frame.paragraphs:
            for run in p.runs: run.font.size = Pt(11)

    # Properties
    set_cell(1,0,"Properties")
    for i,ag in enumerate(AGENTS, start=1):
        set_cell(1,i,"\n".join(props[ag]) or "—")

    # Objectives
    set_cell(2,0,"Objectives")
    for i,ag in enumerate(AGENTS, start=1):
        set_cell(2,i,"\n".join(objectives.get(ag,[])) or "—")

    # States
    set_cell(3,0,"States")
    for i,ag in enumerate(AGENTS, start=1):
        s = st[ag]
        txt = []
        txt.append("Self:")
        txt += s.get("Self",[]) or ["—"]
        txt.append("")
        txt.append("Neighbors:")
        txt += s.get("Neighbors",[]) or ["—"]
        txt.append("")
        txt.append("Relations:")
        txt += s.get("Relations",[]) or ["—"]
        set_cell(3,i,"\n".join(txt))

    # Actions
    set_cell(4,0,"Actions")
    for i,ag in enumerate(AGENTS, start=1):
        lst = acts[ag]
        if not lst:
            set_cell(4,i,"—"); continue
        lines = []
        for a in lst:
            t = a.get("type","")
            ps = a.get("params",[])
            lines.append(f"{t}: {', '.join(ps)}" if ps else t)
        set_cell(4,i,"\n".join(lines))

    prs.save(outpath)

def main():
    os.makedirs(OUT, exist_ok=True)
    props, st, acts, objectives = load_all()
    render_html(props, st, acts, objectives, os.path.join(OUT, "supplychain_table.html"))
    render_pptx(props, st, acts, objectives, os.path.join(OUT, "supplychain_table.pptx"))
    print("✅ Wrote:")
    print("  - outputs/supplychain_table.html")
    print("  - outputs/supplychain_table.pptx")

if __name__ == "__main__":
    main()
